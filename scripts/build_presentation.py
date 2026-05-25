"""Generate the ClearClause sessional presentation as a .pptx.

Run with the project venv:
    /home/shehroz/Documents/genai-course/.venv/bin/python scripts/build_presentation.py

The output lands at ``ClearClause_Presentation.pptx`` in the repo root.

The deck mirrors the LegalTech colour palette of the live Streamlit app
(forest + coral + amber + cream) and is sized for 16:9 widescreen. Each
slide is hand-positioned with python-pptx rather than relying on default
PowerPoint layouts so the typography stays consistent.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# -----------------------------------------------------------------------------
# Palette (matches assets/styles.css)
# -----------------------------------------------------------------------------

FOREST = RGBColor(0x1F, 0x3A, 0x2D)
FOREST_DEEP = RGBColor(0x14, 0x2A, 0x20)
FOREST_BRIGHT = RGBColor(0x2F, 0x5A, 0x40)
MOSS = RGBColor(0x6F, 0x8D, 0x5C)
CORAL = RGBColor(0xC2, 0x41, 0x0C)
CORAL_DEEP = RGBColor(0x9A, 0x1B, 0x0A)
CORAL_SOFT = RGBColor(0xFE, 0xE4, 0xD6)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_SOFT = RGBColor(0xFD, 0xF3, 0xC6)
GREEN = RGBColor(0x15, 0x80, 0x3D)
GREEN_SOFT = RGBColor(0xD6, 0xF0, 0xDF)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_SOFT = RGBColor(0xDB, 0xEA, 0xFE)
INK = RGBColor(0x0F, 0x1D, 0x18)
INK_SOFT = RGBColor(0x24, 0x30, 0x29)
MUTED = RGBColor(0x5A, 0x6A, 0x62)
PAPER = RGBColor(0xF7, 0xF7, 0xF1)
PAPER_SOFT = RGBColor(0xFB, 0xFB, 0xF6)
LINE = RGBColor(0xE3, 0xE6, 0xDD)
LINE_SOFT = RGBColor(0xEE, 0xF0, 0xE8)
CREAM = RGBColor(0xF6, 0xF5, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xB8, 0x8C, 0x3B)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HEADING_FONT = "Cambria"
BODY_FONT = "Calibri"


# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill, line=None, line_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_round(slide, left, top, width, height, fill, line=None, line_w=0.75, radius=0.07):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    try:
        shape.adjustments[0] = radius
    except (IndexError, AttributeError):
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, text, left, top, width, height,
    size=18, bold=False, color=INK, font=BODY_FONT,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_multiline(
    slide, lines, left, top, width, height,
    size=14, color=INK, font=BODY_FONT, bold=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
):
    """Add a multi-paragraph textbox where each entry is one paragraph."""

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, bullets, left, top, width, height,
                size=14, color=INK_SOFT, bullet_color=FOREST_BRIGHT,
                bullet_char="●", line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for idx, text in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        bullet = p.add_run()
        bullet.text = f"{bullet_char}  "
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        bullet.font.name = BODY_FONT
        bullet.font.bold = True
        body = p.add_run()
        body.text = text
        body.font.size = Pt(size)
        body.font.color.rgb = color
        body.font.name = BODY_FONT
    return tb


def slide_background(slide, fill=PAPER_SOFT):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=fill)


def slide_footer(slide, page: int, total: int, label: str = "ClearClause · NLP Final Project · Spring 2026"):
    add_rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), fill=PAPER)
    add_text(slide, label, Inches(0.5), Inches(7.22), Inches(8.5), Inches(0.28),
             size=9, color=MUTED, font=BODY_FONT)
    add_text(slide, f"{page} / {total}", Inches(11.2), Inches(7.22), Inches(1.6), Inches(0.28),
             size=9, color=MUTED, font=BODY_FONT, align=PP_ALIGN.RIGHT)


def slide_header(slide, eyebrow: str, title: str, accent_color=FOREST_BRIGHT):
    # accent bar
    add_rect(slide, Inches(0.45), Inches(0.55), Inches(0.08), Inches(0.7), fill=accent_color)
    add_text(slide, eyebrow.upper(), Inches(0.7), Inches(0.5), Inches(8), Inches(0.32),
             size=10, color=MUTED, bold=True, font=BODY_FONT)
    add_text(slide, title, Inches(0.7), Inches(0.78), Inches(12), Inches(0.7),
             size=30, color=INK, bold=True, font=HEADING_FONT)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# -----------------------------------------------------------------------------
# Slides
# -----------------------------------------------------------------------------

def slide_title(prs, page, total):
    s = blank(prs)
    # forest gradient background built with stacked rectangles
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=FOREST_DEEP)
    add_rect(s, 0, 0, SLIDE_W, Inches(4.4), fill=FOREST)
    add_rect(s, 0, 0, SLIDE_W, Inches(2.0), fill=FOREST_BRIGHT)
    # gold accent stripe
    add_rect(s, 0, Inches(2.0), SLIDE_W, Inches(0.06), fill=GOLD)

    # brand mark
    add_round(s, Inches(0.9), Inches(0.9), Inches(0.85), Inches(0.85), fill=CREAM, radius=0.18)
    add_text(s, "CC", Inches(0.9), Inches(0.95), Inches(0.85), Inches(0.85),
             size=28, bold=True, color=FOREST_DEEP, font=HEADING_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "AI Contract Risk Analyzer", Inches(2.0), Inches(0.95), Inches(8.0), Inches(0.45),
             size=14, color=CREAM, font=BODY_FONT, bold=False)
    add_text(s, "ClearClause", Inches(2.0), Inches(1.25), Inches(10), Inches(0.85),
             size=42, bold=True, color=WHITE, font=HEADING_FONT)

    # tagline below the gold bar
    add_text(s, "TURNING DENSE LEGAL TEXT INTO 60-SECOND DECISIONS.",
             Inches(0.9), Inches(2.45), Inches(12), Inches(0.4),
             size=12, color=CREAM, bold=True, font=BODY_FONT)

    # subtitle paragraph
    add_multiline(s, [
        "A hybrid NLP system that segments freelance and professional contracts,",
        "detects 17 categories of risky clauses, scores the document, benchmarks",
        "each clause against a fair-contract standard, and coaches the user",
        "through the specific negotiation moves to make before signing.",
    ], Inches(0.9), Inches(2.95), Inches(11), Inches(2.0),
       size=15, color=CREAM, line_spacing=1.4)

    # group card
    card_left = Inches(0.9)
    card_top = Inches(5.05)
    card_width = Inches(11.5)
    card_height = Inches(1.85)
    add_round(s, card_left, card_top, card_width, card_height, fill=WHITE)
    add_text(s, "PROJECT GROUP", card_left + Inches(0.4), card_top + Inches(0.2),
             Inches(4), Inches(0.3), size=10, bold=True, color=MUTED, font=BODY_FONT)

    # member columns
    add_text(s, "Shehroz Ali", card_left + Inches(0.4), card_top + Inches(0.55),
             Inches(5), Inches(0.5), size=22, bold=True, color=INK, font=HEADING_FONT)
    add_text(s, "Roll No · MSDSF25M012", card_left + Inches(0.4), card_top + Inches(1.05),
             Inches(5), Inches(0.4), size=13, color=MUTED, font=BODY_FONT)

    add_rect(s, card_left + Inches(5.55), card_top + Inches(0.55),
             Emu(8000), Inches(1.0), fill=LINE)

    add_text(s, "Arslan Ahmad", card_left + Inches(5.85), card_top + Inches(0.55),
             Inches(5), Inches(0.5), size=22, bold=True, color=INK, font=HEADING_FONT)
    add_text(s, "Roll No · MSDSF25M001", card_left + Inches(5.85), card_top + Inches(1.05),
             Inches(5), Inches(0.4), size=13, color=MUTED, font=BODY_FONT)

    add_text(s, "MS Data Science · Natural Language Processing · Spring 2026 Final Project",
             card_left + Inches(0.4), card_top + Inches(1.45),
             Inches(11), Inches(0.3), size=11, color=FOREST_BRIGHT, bold=True, font=BODY_FONT)

    # bottom note
    add_text(s, "github.com/ShehrozRafaqat/ClearClause-NLP-Project",
             Inches(0.0), Inches(7.18), SLIDE_W, Inches(0.32),
             size=10, color=CREAM, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 0, Inches(7.18), SLIDE_W, Inches(0.32), fill=FOREST_DEEP)
    add_text(s, "github.com/ShehrozRafaqat/ClearClause-NLP-Project",
             Inches(0.0), Inches(7.18), SLIDE_W, Inches(0.32),
             size=10, color=CREAM, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def stat_card(slide, left, top, width, height, value, label, color):
    add_round(slide, left, top, width, height, fill=WHITE, line=LINE)
    add_rect(slide, left, top, Inches(0.08), height, fill=color)
    add_text(slide, value, left + Inches(0.3), top + Inches(0.3), width - Inches(0.4), Inches(1.1),
             size=44, bold=True, color=color, font=HEADING_FONT)
    add_text(slide, label, left + Inches(0.3), top + Inches(1.45), width - Inches(0.4), height - Inches(1.5),
             size=12, color=INK_SOFT, font=BODY_FONT)


def slide_problem(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "The problem", "Contracts are written for lawyers, not freelancers.", accent_color=CORAL)

    add_text(s,
             "Most freelancers and small businesses sign multi-page agreements without "
             "fully understanding what they're agreeing to. Risky clauses around IP, "
             "termination, liability and exclusivity routinely cost the smaller party "
             "more than the project fee itself.",
             Inches(0.7), Inches(1.65), Inches(12), Inches(1.4),
             size=15, color=INK_SOFT, font=BODY_FONT)

    # three stat cards
    cards = [
        ("~4,800", "words in the typical freelance contract analysed by ClearClause", CORAL),
        ("17", "clause categories tracked, from IP and indemnity to auto-renewal", AMBER),
        ("96/100", "max risk score the highest-risk demo earns before suggested edits", FOREST_BRIGHT),
    ]
    card_w = Inches(3.95)
    card_h = Inches(2.4)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    for i, (value, label, color) in enumerate(cards):
        left = start_x + i * (card_w + gap)
        stat_card(s, left, Inches(3.4), card_w, card_h, value, label, color)

    # quote band
    add_round(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.9),
              fill=CREAM, line=LINE)
    add_text(s,
             "“Even when the words are clear, the consequences are not. A small wording change in section 9 "
             "is often the difference between getting paid in 15 days and waiting 90.”",
             Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.8),
             size=13, italic=True, color=INK_SOFT, font=BODY_FONT,
             anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(s, page, total)


def pillar_card(slide, left, top, width, height, title, body, color):
    add_round(slide, left, top, width, height, fill=WHITE, line=LINE)
    add_rect(slide, left, top, width, Inches(0.18), fill=color)
    add_text(slide, title, left + Inches(0.3), top + Inches(0.35), width - Inches(0.5), Inches(0.55),
             size=18, bold=True, color=INK, font=HEADING_FONT)
    add_multiline(slide, body, left + Inches(0.3), top + Inches(0.95),
                  width - Inches(0.5), height - Inches(1.05),
                  size=12, color=INK_SOFT, line_spacing=1.35)


def slide_solution(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Our solution", "ClearClause — a hybrid-NLP contract review copilot.")

    add_text(s,
             "Upload a contract; ClearClause segments, classifies, scores, explains, benchmarks, "
             "answers questions, and exports — entirely offline. An optional Groq layer polishes wording.",
             Inches(0.7), Inches(1.65), Inches(12), Inches(0.9),
             size=14, color=INK_SOFT, font=BODY_FONT)

    pillars = [
        ("Detect",
         ["17 clause categories", "Regex catalog + TF-IDF semantic similarity", "Every match is traceable to evidence"],
         FOREST_BRIGHT),
        ("Score",
         ["0–100 risk score with verdict card", "Severity adjustment via red / green flags", "Exponential curve avoids saturation"],
         CORAL),
        ("Coach",
         ["Priority-sorted negotiation checklist", "Fair-contract benchmark per clause", "Suggested discovery questions"],
         AMBER),
        ("Export",
         ["Presentation-quality HTML report", "Clause CSV + Markdown pack", "Document-grounded Q&A chat"],
         GREEN),
    ]
    card_w = Inches(2.92)
    card_h = Inches(3.6)
    gap = Inches(0.22)
    total_w = card_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) / 2
    for i, (title, body, color) in enumerate(pillars):
        left = start_x + i * (card_w + gap)
        pillar_card(s, left, Inches(2.75), card_w, card_h, title, body, color)

    slide_footer(s, page, total)


def pipeline_step(slide, left, top, width, height, idx, title, subtitle, color):
    add_round(slide, left, top, width, height, fill=WHITE, line=LINE)
    add_round(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.55), Inches(0.55),
              fill=color, radius=0.4)
    add_text(slide, str(idx), left + Inches(0.3), top + Inches(0.32),
             Inches(0.55), Inches(0.55), size=15, bold=True, color=WHITE,
             font=HEADING_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, left + Inches(1.0), top + Inches(0.3),
             width - Inches(1.2), Inches(0.55),
             size=15, bold=True, color=INK, font=HEADING_FONT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, subtitle, left + Inches(0.3), top + Inches(1.0),
             width - Inches(0.5), height - Inches(1.05),
             size=11, color=INK_SOFT, font=BODY_FONT)


def slide_pipeline(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Architecture", "The NLP pipeline at a glance.")

    add_text(s,
             "Each step has a single responsibility and a clean dataclass interface, "
             "so the whole pipeline is explainable end-to-end.",
             Inches(0.7), Inches(1.6), Inches(12), Inches(0.6),
             size=13, color=INK_SOFT, font=BODY_FONT)

    # 6 step cards arranged in 2 rows × 3 cols
    steps = [
        (1, "Parse", "TXT / PDF / DOCX → cleaned text + paragraphs.", FOREST_BRIGHT),
        (2, "Segment", "Block-level section splitter (\\Z-anchored regex).", FOREST_BRIGHT),
        (3, "Extract", "Regex patterns + TF-IDF semantic similarity + flag counts → confidence.", CORAL),
        (4, "Score", "Confidence-weighted weights, severity bonus/discount, exponential curve.", CORAL),
        (5, "Benchmark", "Compare each clause to fair-standard text; flag missing core clauses.", AMBER),
        (6, "Coach + Q&A", "Priority checklist, suggested questions, TF-IDF retrieval Q&A.", GREEN),
    ]
    card_w = Inches(4.1)
    card_h = Inches(2.05)
    h_gap = Inches(0.18)
    v_gap = Inches(0.22)
    grid_w = card_w * 3 + h_gap * 2
    start_x = (SLIDE_W - grid_w) / 2
    for i, (idx, title, sub, color) in enumerate(steps):
        row = i // 3
        col = i % 3
        left = start_x + col * (card_w + h_gap)
        top = Inches(2.3) + row * (card_h + v_gap)
        pipeline_step(s, left, top, card_w, card_h, idx, title, sub, color)

    # arrow band
    add_round(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.55),
              fill=CREAM, line=LINE)
    add_text(s,
             "ParsedDocument  →  Section[]  →  ClauseFinding[]  →  RiskProfile  →  AnalysisResult  →  HTML / CSV / Markdown",
             Inches(0.9), Inches(6.7), Inches(11.6), Inches(0.45),
             size=12, bold=True, color=FOREST_DEEP, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)

    slide_footer(s, page, total)


def chip(slide, left, top, width, height, label, fill, text_color):
    add_round(slide, left, top, width, height, fill=fill, radius=0.45)
    add_text(slide, label, left, top, width, height, size=11, bold=True,
             color=text_color, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_catalog(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "The clause catalog", "17 clause categories grouped into 5 risk dimensions.")

    groups = [
        ("Money & Liability", CORAL, CORAL_SOFT, [
            "Payment Terms", "Unlimited Liability", "Indemnification", "Limitation of Liability",
        ]),
        ("Ownership & IP", AMBER, AMBER_SOFT, [
            "IP Ownership / Work-for-Hire", "Confidentiality / NDA", "Data Privacy / Security",
        ]),
        ("Termination & Renewal", FOREST_BRIGHT, GREEN_SOFT, [
            "Unilateral Termination", "Automatic Renewal",
        ]),
        ("Work Freedom", GOLD, AMBER_SOFT, [
            "Non-Compete", "Exclusivity", "Non-Solicitation",
        ]),
        ("Disputes", BLUE, BLUE_SOFT, [
            "Governing Law & Jurisdiction", "Arbitration / Court Waiver",
        ]),
        ("Operations", MOSS, GREEN_SOFT, [
            "Scope, Acceptance & Revisions", "Assignment / Change of Control", "Effective Date & Term",
        ]),
    ]

    # grid layout: 2 cols × 3 rows
    col_w = Inches(6.05)
    row_h = Inches(1.65)
    h_gap = Inches(0.25)
    v_gap = Inches(0.15)
    start_x = (SLIDE_W - col_w * 2 - h_gap) / 2
    start_y = Inches(1.85)

    for i, (group_name, accent, soft, items) in enumerate(groups):
        col = i % 2
        row = i // 2
        left = start_x + col * (col_w + h_gap)
        top = start_y + row * (row_h + v_gap)
        # group card
        add_round(s, left, top, col_w, row_h, fill=WHITE, line=LINE)
        add_rect(s, left, top, Inches(0.1), row_h, fill=accent)
        add_text(s, group_name, left + Inches(0.3), top + Inches(0.18),
                 col_w - Inches(0.5), Inches(0.4),
                 size=14, bold=True, color=INK, font=HEADING_FONT)
        # chips
        chip_y = top + Inches(0.65)
        chip_x = left + Inches(0.3)
        chip_h = Inches(0.4)
        running_x = chip_x
        for item in items:
            estimated = max(Inches(1.2), Inches(0.05 * len(item) + 0.7))
            if (running_x - left) + estimated > col_w - Inches(0.2):
                running_x = chip_x
                chip_y += Inches(0.5)
            chip(s, running_x, chip_y, estimated, chip_h, item, soft, accent)
            running_x += estimated + Inches(0.1)

    slide_footer(s, page, total)


def slide_hybrid_extraction(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "How detection works", "Three signals combine into one explainable confidence score.")

    # left column: signals
    sig_x = Inches(0.7)
    sig_y = Inches(1.75)
    sig_w = Inches(5.9)

    add_round(s, sig_x, sig_y, sig_w, Inches(1.7), fill=WHITE, line=LINE)
    add_rect(s, sig_x, sig_y, Inches(0.1), Inches(1.7), fill=FOREST_BRIGHT)
    add_text(s, "1 · Regex catalog match", sig_x + Inches(0.3), sig_y + Inches(0.15),
             sig_w - Inches(0.5), Inches(0.4), size=14, bold=True, color=INK, font=HEADING_FONT)
    add_multiline(s, [
        "Hand-curated patterns and red/green flag terms per clause.",
        "Captures the explicit, traceable signal.",
        "Example: \\birrevocably assigns?\\b → IP Ownership.",
    ], sig_x + Inches(0.3), sig_y + Inches(0.55), sig_w - Inches(0.5), Inches(1.1),
       size=11, color=INK_SOFT, line_spacing=1.3)

    add_round(s, sig_x, sig_y + Inches(1.85), sig_w, Inches(1.7), fill=WHITE, line=LINE)
    add_rect(s, sig_x, sig_y + Inches(1.85), Inches(0.1), Inches(1.7), fill=CORAL)
    add_text(s, "2 · TF-IDF semantic similarity",
             sig_x + Inches(0.3), sig_y + Inches(2.0),
             sig_w - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK, font=HEADING_FONT)
    add_multiline(s, [
        "Cosine similarity between each section and the rule's prototype text.",
        "Catches paraphrased or unusual wording the regex misses.",
        "Adds a continuous signal alongside binary pattern matches.",
    ], sig_x + Inches(0.3), sig_y + Inches(2.4), sig_w - Inches(0.5), Inches(1.1),
       size=11, color=INK_SOFT, line_spacing=1.3)

    add_round(s, sig_x, sig_y + Inches(3.7), sig_w, Inches(1.7), fill=WHITE, line=LINE)
    add_rect(s, sig_x, sig_y + Inches(3.7), Inches(0.1), Inches(1.7), fill=AMBER)
    add_text(s, "3 · Red / green flag counters",
             sig_x + Inches(0.3), sig_y + Inches(3.85),
             sig_w - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK, font=HEADING_FONT)
    add_multiline(s, [
        "Red flags (e.g. 'unlimited', 'sole discretion') push severity up.",
        "Green/balancing flags (e.g. 'fees paid', 'mutual') pull severity down.",
        "Drives the per-clause severity adjustment.",
    ], sig_x + Inches(0.3), sig_y + Inches(4.25), sig_w - Inches(0.5), Inches(1.1),
       size=11, color=INK_SOFT, line_spacing=1.3)

    # right column: worked example
    ex_x = Inches(7.0)
    ex_y = Inches(1.75)
    ex_w = Inches(5.65)
    ex_h = Inches(5.4)
    add_round(s, ex_x, ex_y, ex_w, ex_h, fill=FOREST_DEEP)
    add_text(s, "WORKED EXAMPLE", ex_x + Inches(0.3), ex_y + Inches(0.25),
             ex_w - Inches(0.5), Inches(0.3), size=10, bold=True, color=CREAM,
             font=BODY_FONT)
    add_text(s, "IP Ownership clause in the high-risk freelance demo",
             ex_x + Inches(0.3), ex_y + Inches(0.55),
             ex_w - Inches(0.5), Inches(0.45),
             size=15, bold=True, color=WHITE, font=HEADING_FONT)

    # snippet card
    snip_y = ex_y + Inches(1.1)
    add_round(s, ex_x + Inches(0.3), snip_y, ex_w - Inches(0.6), Inches(1.5),
              fill=PAPER, line=None)
    add_text(s,
             "\"To the extent any work product is not work-for-hire, Contractor "
             "irrevocably assigns all right, title, and interest in such intellectual "
             "property to Client. Client shall exclusively own all deliverables.\"",
             ex_x + Inches(0.45), snip_y + Inches(0.15),
             ex_w - Inches(0.9), Inches(1.2),
             size=11, italic=True, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)

    # signals breakdown
    metrics_y = snip_y + Inches(1.7)
    rows = [
        ("Regex pattern matches", "intellectual property · irrevocably assigns · all right, title, interest"),
        ("Red flag terms", "irrevocably · exclusively owned · all right"),
        ("Green / balancing terms", "(none)"),
        ("TF-IDF similarity to prototype", "0.42"),
        ("Final confidence", "0.98"),
        ("Adjusted severity", "HIGH (no demotion — no balancing language)"),
    ]
    for i, (label, value) in enumerate(rows):
        row_y = metrics_y + i * Inches(0.36)
        add_text(s, label, ex_x + Inches(0.3), row_y,
                 Inches(2.6), Inches(0.32), size=10, bold=True, color=CREAM, font=BODY_FONT)
        add_text(s, value, ex_x + Inches(2.95), row_y,
                 ex_w - Inches(3.25), Inches(0.32), size=10, color=CREAM, font=BODY_FONT)

    slide_footer(s, page, total)


def slide_scoring(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Risk score", "Transparent formula. No saturation at the top.")

    # left side: formula + verdict thresholds
    left_x = Inches(0.7)
    left_w = Inches(6.0)
    add_round(s, left_x, Inches(1.7), left_w, Inches(3.3),
              fill=WHITE, line=LINE)
    add_text(s, "Formula", left_x + Inches(0.3), Inches(1.85),
             left_w - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK, font=HEADING_FONT)
    add_multiline(s, [
        "contribution = (weight + red_bonus − green_discount) × conf_factor",
        "raw_total = Σ contributions + 3.5·high_count + 1.25·medium_count",
        "if raw_total ≤ 60:    score = round(raw_total)",
        "else:                  score = round(60 + 36·(1 − exp(−(raw_total − 60) / 160)))",
        "                       ⇒ asymptote ≈ 96, ceiling = 100",
    ], left_x + Inches(0.3), Inches(2.3), left_w - Inches(0.5), Inches(2.65),
       size=11, color=INK_SOFT, line_spacing=1.5, font="Consolas")

    # verdict thresholds card
    add_round(s, left_x, Inches(5.15), left_w, Inches(1.9),
              fill=WHITE, line=LINE)
    add_text(s, "Verdict thresholds", left_x + Inches(0.3), Inches(5.25),
             left_w - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK, font=HEADING_FONT)
    verdicts = [
        ("≥ 78", "Critical Risk", "Do not sign in current form.", CORAL_DEEP),
        ("60 – 77", "High Risk", "Revisions required before signing.", CORAL),
        ("30 – 59", "Review Carefully", "Manageable; clarify wording.", AMBER),
        ("< 30", "Low Risk", "Final read-through recommended.", GREEN),
    ]
    for i, (band, label, body, color) in enumerate(verdicts):
        row_y = Inches(5.7) + i * Inches(0.32)
        add_round(s, left_x + Inches(0.3), row_y, Inches(1.0), Inches(0.32),
                  fill=color, radius=0.4)
        add_text(s, band, left_x + Inches(0.3), row_y,
                 Inches(1.0), Inches(0.32),
                 size=10, bold=True, color=WHITE, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, label, left_x + Inches(1.45), row_y,
                 Inches(1.6), Inches(0.32),
                 size=11, bold=True, color=color, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, body, left_x + Inches(3.0), row_y,
                 left_w - Inches(3.3), Inches(0.32),
                 size=10, color=INK_SOFT, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # right side: before/after chart visualization
    right_x = Inches(7.0)
    right_w = Inches(5.65)
    add_round(s, right_x, Inches(1.7), right_w, Inches(5.35),
              fill=PAPER, line=LINE)
    add_text(s, "Score saturation, before vs. after",
             right_x + Inches(0.3), Inches(1.85),
             right_w - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK, font=HEADING_FONT)
    add_text(s, "Before the curve, two very different bad contracts both hit 100. After the curve they spread out — the dashboard finally has range.",
             right_x + Inches(0.3), Inches(2.3),
             right_w - Inches(0.5), Inches(1.0),
             size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)

    # before/after bar pairs for two demos
    pair_y = Inches(3.5)
    bar_max_w = Inches(3.8)
    bar_left = right_x + Inches(1.3)

    # row 1: High-risk freelance
    add_text(s, "Freelance", right_x + Inches(0.3), pair_y,
             Inches(1.0), Inches(0.3), size=10, color=INK, bold=True, font=BODY_FONT)
    add_rect(s, bar_left, pair_y + Inches(0.04), bar_max_w * (100/100), Inches(0.25),
             fill=CORAL_DEEP)
    add_text(s, "before · 100",
             bar_left + Inches(0.1), pair_y + Inches(0.04),
             bar_max_w, Inches(0.25),
             size=9, bold=True, color=WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, bar_left, pair_y + Inches(0.4), bar_max_w * (90/100), Inches(0.25),
             fill=CORAL)
    add_text(s, "after · 90",
             bar_left + Inches(0.1), pair_y + Inches(0.4),
             bar_max_w, Inches(0.25),
             size=9, bold=True, color=WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # row 2: SaaS subscription
    pair_y = Inches(4.4)
    add_text(s, "SaaS", right_x + Inches(0.3), pair_y,
             Inches(1.0), Inches(0.3), size=10, color=INK, bold=True, font=BODY_FONT)
    add_rect(s, bar_left, pair_y + Inches(0.04), bar_max_w * (100/100), Inches(0.25),
             fill=CORAL_DEEP)
    add_text(s, "before · 100",
             bar_left + Inches(0.1), pair_y + Inches(0.04),
             bar_max_w, Inches(0.25),
             size=9, bold=True, color=WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, bar_left, pair_y + Inches(0.4), bar_max_w * (85/100), Inches(0.25),
             fill=CORAL)
    add_text(s, "after · 85",
             bar_left + Inches(0.1), pair_y + Inches(0.4),
             bar_max_w, Inches(0.25),
             size=9, bold=True, color=WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # row 3: held-out
    pair_y = Inches(5.3)
    add_text(s, "Influencer", right_x + Inches(0.3), pair_y,
             Inches(1.0), Inches(0.3), size=10, color=INK, bold=True, font=BODY_FONT)
    add_rect(s, bar_left, pair_y + Inches(0.04), bar_max_w * (100/100), Inches(0.25),
             fill=CORAL_DEEP)
    add_text(s, "before · 100",
             bar_left + Inches(0.1), pair_y + Inches(0.04),
             bar_max_w, Inches(0.25),
             size=9, bold=True, color=WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, bar_left, pair_y + Inches(0.4), bar_max_w * (80/100), Inches(0.25),
             fill=AMBER)
    add_text(s, "after · 80",
             bar_left + Inches(0.1), pair_y + Inches(0.4),
             bar_max_w, Inches(0.25),
             size=9, bold=True, color=WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Spread now 90 → 85 → 80 → 53 → 25 across all 5 demos.",
             right_x + Inches(0.3), Inches(6.55),
             right_w - Inches(0.5), Inches(0.35),
             size=10, italic=True, color=MUTED, font=BODY_FONT)

    slide_footer(s, page, total)


def slide_negotiation(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Negotiation coach", "From detection to action — what to ask for.")

    # left column: checklist card mock
    cl_x = Inches(0.7)
    cl_w = Inches(6.0)
    cl_h = Inches(5.1)
    add_round(s, cl_x, Inches(1.75), cl_w, cl_h, fill=WHITE, line=LINE)
    add_rect(s, cl_x, Inches(1.75), Inches(0.1), cl_h, fill=CORAL)

    add_text(s, "#1 · Indemnification", cl_x + Inches(0.3), Inches(1.95),
             cl_w - Inches(0.5), Inches(0.5),
             size=18, bold=True, color=INK, font=HEADING_FONT)
    # severity badge
    add_round(s, cl_x + cl_w - Inches(1.2), Inches(2.05), Inches(0.85), Inches(0.32),
              fill=CORAL_SOFT, radius=0.45)
    add_text(s, "HIGH · 98%", cl_x + cl_w - Inches(1.2), Inches(2.05),
             Inches(0.85), Inches(0.32), size=9, bold=True, color=CORAL,
             font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Legal-defense costs can be much higher than the project value.",
             cl_x + Inches(0.3), Inches(2.55),
             cl_w - Inches(0.5), Inches(0.45),
             size=12, italic=True, color=MUTED, font=BODY_FONT)

    actions = [
        "Make indemnity mutual: each party covers claims it actually caused",
        "Limit the trigger to 'to the extent caused by' your own breach, negligence, or IP claim",
        "Require the client to control defense and let you approve any settlement",
    ]
    for i, action in enumerate(actions):
        y = Inches(3.15) + i * Inches(0.55)
        add_round(s, cl_x + Inches(0.3), y + Inches(0.1), Inches(0.22), Inches(0.22),
                  fill=FOREST_BRIGHT, radius=0.45)
        add_text(s, "✓", cl_x + Inches(0.3), y + Inches(0.07),
                 Inches(0.22), Inches(0.22), size=10, bold=True, color=WHITE,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, action, cl_x + Inches(0.7), y,
                 cl_w - Inches(1.0), Inches(0.4),
                 size=12, color=INK_SOFT, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # right column: fair benchmark table mock
    bx = Inches(7.0)
    bw = Inches(5.65)
    bh = Inches(5.1)
    add_round(s, bx, Inches(1.75), bw, bh, fill=WHITE, line=LINE)
    add_text(s, "Fair-contract benchmark", bx + Inches(0.3), Inches(1.95),
             bw - Inches(0.5), Inches(0.4), size=14, bold=True, color=INK, font=HEADING_FONT)
    add_text(s, "Each detected clause is compared to a fair-standard version. Missing standard clauses are flagged too.",
             bx + Inches(0.3), Inches(2.35),
             bw - Inches(0.5), Inches(0.7),
             size=11, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    # header row
    add_rect(s, bx + Inches(0.3), Inches(3.2), bw - Inches(0.6), Inches(0.36), fill=PAPER)
    add_text(s, "Clause", bx + Inches(0.45), Inches(3.21),
             Inches(2), Inches(0.36), size=10, bold=True, color=INK,
             font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Status", bx + Inches(2.6), Inches(3.21),
             Inches(1.2), Inches(0.36), size=10, bold=True, color=INK,
             font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Detected language vs. fair standard", bx + Inches(3.85), Inches(3.21),
             Inches(2), Inches(0.36), size=10, bold=True, color=INK,
             font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("IP Ownership", "high-risk", CORAL_SOFT, CORAL,
         "Irrevocable assignment with no carve-outs."),
        ("Unilateral Termination", "high-risk", CORAL_SOFT, CORAL,
         "Sole-discretion termination, no kill fee."),
        ("Liability Cap", "concerning", AMBER_SOFT, AMBER,
         "One-sided cap protects only the client."),
        ("Confidentiality", "fair", GREEN_SOFT, GREEN,
         "Mutual, 3-year term, standard exceptions."),
        ("Non-Solicitation", "missing", BLUE_SOFT, BLUE,
         "Standard clause not present — ask why."),
    ]
    row_h = Inches(0.34)
    row_gap = Inches(0.05)
    for i, (clause, status, soft, color, text) in enumerate(rows):
        y = Inches(3.6) + i * (row_h + row_gap)
        add_text(s, clause, bx + Inches(0.3), y,
                 Inches(2.0), row_h, size=10, bold=True, color=INK, font=BODY_FONT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, bx + Inches(2.4), y + Inches(0.05), Inches(1.15), row_h - Inches(0.1),
                  fill=soft, radius=0.45)
        add_text(s, status, bx + Inches(2.4), y + Inches(0.05),
                 Inches(1.15), row_h - Inches(0.1), size=8, bold=True, color=color,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, text, bx + Inches(3.65), y,
                 bw - Inches(3.95), row_h,
                 size=9, color=INK_SOFT, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(s, page, total)


def slide_redline(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Inline redline view", "Full contract with risky spans highlighted in place.")

    intro_x = Inches(0.7)
    intro_y = Inches(1.7)
    add_text(s,
             "The Streamlit redline view stitches all detected spans back into the source text so a "
             "reviewer can see at a glance which paragraphs need attention — hover any highlight for the "
             "clause name.",
             intro_x, intro_y, Inches(12), Inches(0.85),
             size=13, color=INK_SOFT, font=BODY_FONT, line_spacing=1.35)

    # legend chips
    chip(s, Inches(0.7), Inches(2.75), Inches(1.2), Inches(0.35),
         "HIGH · 7", CORAL_SOFT, CORAL)
    chip(s, Inches(2.0), Inches(2.75), Inches(1.4), Inches(0.35),
         "MEDIUM · 7", AMBER_SOFT, AMBER)
    chip(s, Inches(3.5), Inches(2.75), Inches(1.2), Inches(0.35),
         "LOW · 2", GREEN_SOFT, GREEN)
    add_text(s, "(from the high-risk freelance demo)",
             Inches(4.85), Inches(2.75), Inches(5), Inches(0.35),
             size=10, italic=True, color=MUTED, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # mock contract panel
    panel_x = Inches(0.7)
    panel_y = Inches(3.25)
    panel_w = Inches(12)
    panel_h = Inches(3.75)
    add_round(s, panel_x, panel_y, panel_w, panel_h, fill=WHITE, line=LINE)
    add_round(s, panel_x + Inches(0.2), panel_y + Inches(0.2),
              panel_w - Inches(0.4), panel_h - Inches(0.4),
              fill=PAPER, line=None)

    # mock highlighted lines (use rectangles with colored backgrounds + text)
    line_y = panel_y + Inches(0.45)
    line_h = Inches(0.35)

    def line(text, fill, color, top, height=line_h):
        if fill is not None:
            add_rect(s, panel_x + Inches(0.4), top + Inches(0.04), panel_w - Inches(0.8),
                     height - Inches(0.08), fill=fill)
            # underline
            add_rect(s, panel_x + Inches(0.4), top + height - Inches(0.06),
                     panel_w - Inches(0.8), Emu(20000), fill=color)
        add_text(s, text, panel_x + Inches(0.5), top,
                 panel_w - Inches(1.0), height,
                 size=11, color=INK if fill else INK_SOFT, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)

    line("3. Intellectual Property", None, None, line_y)
    line("Contractor irrevocably assigns all right, title, and interest to Client.",
         CORAL_SOFT, CORAL, line_y + Inches(0.35))
    line("5. Non-Compete", None, None, line_y + Inches(0.85))
    line("Contractor shall not compete with Client for 24 months after termination.",
         CORAL_SOFT, CORAL, line_y + Inches(1.2))
    line("8. Termination", None, None, line_y + Inches(1.7))
    line("Client may terminate this Agreement at its convenience upon 7 days notice.",
         CORAL_SOFT, CORAL, line_y + Inches(2.05))
    line("12. Governing Law", None, None, line_y + Inches(2.55))
    line("Governed by the laws of Delaware; exclusive jurisdiction in New Castle County.",
         AMBER_SOFT, AMBER, line_y + Inches(2.9))

    slide_footer(s, page, total)


def slide_qa(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Ask the document", "Grounded retrieval Q&A — no hallucinations.")

    # left side flow
    lx = Inches(0.7)
    lw = Inches(6.0)
    add_text(s, "Retrieval flow",
             lx, Inches(1.75), lw, Inches(0.4),
             size=15, bold=True, color=INK, font=HEADING_FONT)

    steps = [
        ("1", "Chunk contract", "Split into ~950-char overlapping paragraphs.", FOREST_BRIGHT),
        ("2", "Match clause intent", "Keyword map for 17 topics → shortcut to detected clause.", AMBER),
        ("3", "TF-IDF retrieve", "Cosine top-k chunks for the user's question.", CORAL),
        ("4", "Respond", "Use the clause's plain-English answer + cite contract evidence.", GREEN),
    ]
    for i, (idx, title, body, color) in enumerate(steps):
        y = Inches(2.25) + i * Inches(1.05)
        add_round(s, lx, y, lw, Inches(0.95), fill=WHITE, line=LINE)
        add_round(s, lx + Inches(0.25), y + Inches(0.18), Inches(0.5), Inches(0.5),
                  fill=color, radius=0.4)
        add_text(s, idx, lx + Inches(0.25), y + Inches(0.18),
                 Inches(0.5), Inches(0.5), size=14, bold=True, color=WHITE,
                 font=HEADING_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, title, lx + Inches(0.9), y + Inches(0.12),
                 lw - Inches(1.2), Inches(0.4),
                 size=13, bold=True, color=INK, font=HEADING_FONT)
        add_text(s, body, lx + Inches(0.9), y + Inches(0.5),
                 lw - Inches(1.2), Inches(0.45),
                 size=11, color=INK_SOFT, font=BODY_FONT)

    # right side: example dialogue
    rx = Inches(7.0)
    rw = Inches(5.65)
    rh = Inches(5.4)
    add_round(s, rx, Inches(1.75), rw, rh, fill=WHITE, line=LINE)
    add_text(s, "Example exchange",
             rx + Inches(0.3), Inches(1.9), rw - Inches(0.5), Inches(0.4),
             size=15, bold=True, color=INK, font=HEADING_FONT)

    # user bubble
    add_round(s, rx + Inches(1.0), Inches(2.4), rw - Inches(1.3), Inches(0.65),
              fill=FOREST_BRIGHT)
    add_text(s, "Who owns the final source code?",
             rx + Inches(1.15), Inches(2.4),
             rw - Inches(1.55), Inches(0.65),
             size=12, color=WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # assistant bubble
    add_round(s, rx + Inches(0.3), Inches(3.2),
              rw - Inches(1.3), Inches(3.6),
              fill=PAPER, line=LINE)
    add_text(s, "Detected: IP Ownership (HIGH · 98% confidence)",
             rx + Inches(0.5), Inches(3.3),
             rw - Inches(1.6), Inches(0.35),
             size=10, bold=True, color=CORAL, font=BODY_FONT)
    add_multiline(s, [
        "Plain English: The contract assigns all IP, source code, and deliverables to the client. You will not own the finished work.",
        "",
        "Negotiation note: Carve out pre-existing tools, templates, and generic libraries; reserve portfolio display rights.",
        "",
        "Evidence: \"Contractor irrevocably assigns all right, title, and interest in such intellectual property to Client. Client shall exclusively own all deliverables.\"",
    ], rx + Inches(0.5), Inches(3.7),
       rw - Inches(1.6), Inches(3.0),
       size=10, color=INK_SOFT, line_spacing=1.35)

    slide_footer(s, page, total)


def slide_dashboard(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Executive dashboard", "Score, verdict, radar, KPI strip, and ranked red flags.")

    # KPI strip
    strip_y = Inches(1.7)
    kpi_data = [
        ("90/100", "Risk score", CORAL_DEEP),
        ("7", "High-risk clauses", CORAL),
        ("7", "Medium-risk clauses", AMBER),
        ("3", "Low + Info clauses", GREEN),
        ("89%", "Avg confidence", FOREST_BRIGHT),
    ]
    kpi_w = Inches(2.4)
    kpi_gap = Inches(0.15)
    total_kpi_w = kpi_w * 5 + kpi_gap * 4
    start_x = (SLIDE_W - total_kpi_w) / 2
    for i, (val, label, color) in enumerate(kpi_data):
        kx = start_x + i * (kpi_w + kpi_gap)
        add_round(s, kx, strip_y, kpi_w, Inches(1.05), fill=WHITE, line=LINE)
        add_text(s, label, kx + Inches(0.2), strip_y + Inches(0.1),
                 kpi_w - Inches(0.3), Inches(0.3),
                 size=9, bold=True, color=MUTED, font=BODY_FONT)
        add_text(s, val, kx + Inches(0.2), strip_y + Inches(0.38),
                 kpi_w - Inches(0.3), Inches(0.6),
                 size=22, bold=True, color=color, font=HEADING_FONT)

    # verdict card (full width)
    v_y = Inches(2.95)
    add_round(s, Inches(0.7), v_y, Inches(12), Inches(1.05), fill=WHITE, line=LINE)
    add_rect(s, Inches(0.7), v_y, Inches(0.12), Inches(1.05), fill=CORAL_DEEP)
    add_text(s, "CLEARCLAUSE VERDICT", Inches(0.95), v_y + Inches(0.1),
             Inches(4), Inches(0.3), size=10, bold=True, color=MUTED, font=BODY_FONT)
    add_text(s, "Critical Risk — Do not sign in the current form.",
             Inches(0.95), v_y + Inches(0.38),
             Inches(11), Inches(0.4),
             size=16, bold=True, color=INK, font=HEADING_FONT)
    add_text(s, "Concentrates several high-severity clauses against one party. Push for revisions on the priority red flags before signing.",
             Inches(0.95), v_y + Inches(0.7),
             Inches(11), Inches(0.32),
             size=10, color=MUTED, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # Three small panels: gauge mock, radar mock, top flags
    panel_y = Inches(4.2)
    panel_h = Inches(2.85)
    panel_w = Inches(3.9)
    panel_gap = Inches(0.15)
    total_panel_w = panel_w * 3 + panel_gap * 2
    panel_start_x = (SLIDE_W - total_panel_w) / 2

    # 1. risk gauge mock
    p1_x = panel_start_x
    add_round(s, p1_x, panel_y, panel_w, panel_h, fill=WHITE, line=LINE)
    add_text(s, "Composite risk score", p1_x + Inches(0.2), panel_y + Inches(0.15),
             panel_w - Inches(0.4), Inches(0.3),
             size=11, bold=True, color=MUTED, font=BODY_FONT)
    # semicircle gauge: use a pie shape would be complex, use stacked rectangles for arc-like look
    gauge_cx = p1_x + panel_w / 2
    gauge_cy = panel_y + Inches(1.5)
    # background arc segments
    segments = [
        (0, 30, GREEN_SOFT),
        (30, 60, AMBER_SOFT),
        (60, 78, CORAL_SOFT),
        (78, 100, RGBColor(0xf8, 0xc8, 0xb9)),
    ]
    seg_w = Inches(2.8)
    seg_h = Inches(0.4)
    seg_y = panel_y + Inches(1.45)
    seg_left = p1_x + Inches(0.55)
    for low, high, color in segments:
        width = seg_w * ((high - low) / 100)
        offset = seg_w * (low / 100)
        add_rect(s, seg_left + offset, seg_y, width, seg_h, fill=color)
    # indicator
    add_rect(s, seg_left + seg_w * 0.9 - Emu(20000), seg_y - Inches(0.1),
             Emu(40000), seg_h + Inches(0.2), fill=FOREST_DEEP)
    # value
    add_text(s, "90/100", p1_x + Inches(0.2), panel_y + Inches(2.0),
             panel_w - Inches(0.4), Inches(0.5),
             size=24, bold=True, color=CORAL, font=HEADING_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Critical Risk",
             p1_x + Inches(0.2), panel_y + Inches(2.45),
             panel_w - Inches(0.4), Inches(0.3),
             size=10, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # 2. radar mock (just a static visualization with dimension labels)
    p2_x = panel_start_x + panel_w + panel_gap
    add_round(s, p2_x, panel_y, panel_w, panel_h, fill=WHITE, line=LINE)
    add_text(s, "Risk by business area", p2_x + Inches(0.2), panel_y + Inches(0.15),
             panel_w - Inches(0.4), Inches(0.3),
             size=11, bold=True, color=MUTED, font=BODY_FONT)

    # use horizontal bars instead of radar for clarity
    dims = [
        ("Money exposure", 92, CORAL),
        ("Ownership & data", 84, AMBER),
        ("Work freedom", 88, CORAL),
        ("Disputes & enforcement", 70, AMBER),
        ("Operational clarity", 45, MOSS),
    ]
    bar_max_w = Inches(2.4)
    for i, (label, val, color) in enumerate(dims):
        y = panel_y + Inches(0.6) + i * Inches(0.4)
        add_text(s, label, p2_x + Inches(0.2), y,
                 Inches(1.5), Inches(0.3), size=9, color=INK_SOFT, font=BODY_FONT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, p2_x + Inches(1.55), y + Inches(0.08),
                 bar_max_w * (val/100), Inches(0.18), fill=color)
        add_text(s, str(val), p2_x + Inches(3.55), y,
                 Inches(0.3), Inches(0.3), size=9, bold=True, color=INK,
                 font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # 3. top flags
    p3_x = panel_start_x + (panel_w + panel_gap) * 2
    add_round(s, p3_x, panel_y, panel_w, panel_h, fill=WHITE, line=LINE)
    add_text(s, "Priority red flags", p3_x + Inches(0.2), panel_y + Inches(0.15),
             panel_w - Inches(0.4), Inches(0.3),
             size=11, bold=True, color=MUTED, font=BODY_FONT)
    flags = [
        ("Indemnification", "Legal-defense costs can exceed project value."),
        ("IP Ownership", "Loss of source-code and portfolio rights."),
        ("Non-Compete", "Blocks main source of income after contract."),
        ("Unlimited Liability", "Personal exposure beyond fee."),
    ]
    for i, (title, body) in enumerate(flags):
        y = panel_y + Inches(0.55) + i * Inches(0.52)
        add_round(s, p3_x + Inches(0.2), y, panel_w - Inches(0.4), Inches(0.45),
                  fill=PAPER, line=None)
        add_rect(s, p3_x + Inches(0.2), y, Inches(0.06), Inches(0.45), fill=CORAL)
        add_text(s, title, p3_x + Inches(0.35), y,
                 Inches(1.7), Inches(0.45),
                 size=10, bold=True, color=INK, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, body, p3_x + Inches(2.05), y,
                 panel_w - Inches(2.25), Inches(0.45),
                 size=8, color=MUTED, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(s, page, total)


def metric_pill(slide, left, top, width, height, label, value, color):
    add_round(slide, left, top, width, height, fill=WHITE, line=LINE)
    add_rect(slide, left, top, Inches(0.08), height, fill=color)
    add_text(slide, label, left + Inches(0.25), top + Inches(0.12),
             width - Inches(0.4), Inches(0.3),
             size=10, bold=True, color=MUTED, font=BODY_FONT)
    add_text(slide, value, left + Inches(0.25), top + Inches(0.4),
             width - Inches(0.4), height - Inches(0.5),
             size=24, bold=True, color=color, font=HEADING_FONT)


def slide_evaluation(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Evaluation", "Gold (in-domain) + held-out (out-of-domain) = honest metrics.")

    add_text(s,
             "The hand-curated gold set scores 1.00 across the board, but that's the optimistic "
             "number because the catalog was tuned against this contract. The held-out demo is from "
             "a completely different domain (influencer marketing) and reports the realistic number.",
             Inches(0.7), Inches(1.6), Inches(12), Inches(0.95),
             size=12, color=INK_SOFT, font=BODY_FONT, line_spacing=1.35)

    # Two evaluation panels
    panel_w = Inches(5.95)
    panel_h = Inches(3.65)
    panel_gap = Inches(0.2)
    start_x = (SLIDE_W - panel_w * 2 - panel_gap) / 2

    # Gold panel
    gx = start_x
    gy = Inches(2.7)
    add_round(s, gx, gy, panel_w, panel_h, fill=WHITE, line=LINE)
    add_rect(s, gx, gy, panel_w, Inches(0.5), fill=FOREST_BRIGHT)
    add_text(s, "Gold demo · high-risk freelance contract",
             gx + Inches(0.3), gy + Inches(0.05),
             panel_w - Inches(0.6), Inches(0.4),
             size=13, bold=True, color=WHITE, font=HEADING_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "17 expected clause IDs · catalog was tuned against this contract",
             gx + Inches(0.3), gy + Inches(0.65),
             panel_w - Inches(0.6), Inches(0.35),
             size=11, color=MUTED, font=BODY_FONT)

    pill_w = (panel_w - Inches(0.85)) / 3
    pill_y = gy + Inches(1.15)
    pill_h = Inches(1.1)
    metric_pill(s, gx + Inches(0.3), pill_y, pill_w, pill_h, "Precision", "1.00", FOREST_BRIGHT)
    metric_pill(s, gx + Inches(0.3) + pill_w + Inches(0.125), pill_y, pill_w, pill_h, "Recall", "1.00", FOREST_BRIGHT)
    metric_pill(s, gx + Inches(0.3) + (pill_w + Inches(0.125)) * 2, pill_y, pill_w, pill_h, "F1", "1.00", FOREST_BRIGHT)

    add_text(s, "TP: 17  ·  FP: 0  ·  FN: 0",
             gx + Inches(0.3), gy + Inches(2.45),
             panel_w - Inches(0.6), Inches(0.3),
             size=12, bold=True, color=INK, font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(s,
             "Every expected clause type is detected. Confidence is high (avg 0.87), "
             "and severity counts (7 HIGH / 8 MEDIUM) match a sensible legal review.",
             gx + Inches(0.3), gy + Inches(2.8),
             panel_w - Inches(0.6), Inches(0.8),
             size=10, color=INK_SOFT, font=BODY_FONT, line_spacing=1.35)

    # Held-out panel
    hx = gx + panel_w + panel_gap
    add_round(s, hx, gy, panel_w, panel_h, fill=WHITE, line=LINE)
    add_rect(s, hx, gy, panel_w, Inches(0.5), fill=CORAL)
    add_text(s, "Held-out · creator / influencer agreement",
             hx + Inches(0.3), gy + Inches(0.05),
             panel_w - Inches(0.6), Inches(0.4),
             size=13, bold=True, color=WHITE, font=HEADING_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "11 expected clause IDs · catalog was NOT tuned for this domain",
             hx + Inches(0.3), gy + Inches(0.65),
             panel_w - Inches(0.6), Inches(0.35),
             size=11, color=MUTED, font=BODY_FONT)

    metric_pill(s, hx + Inches(0.3), pill_y, pill_w, pill_h, "Precision", "0.92", CORAL)
    metric_pill(s, hx + Inches(0.3) + pill_w + Inches(0.125), pill_y, pill_w, pill_h, "Recall", "1.00", FOREST_BRIGHT)
    metric_pill(s, hx + Inches(0.3) + (pill_w + Inches(0.125)) * 2, pill_y, pill_w, pill_h, "F1", "0.96", AMBER)

    add_text(s, "TP: 11  ·  FP: 1 (non_compete)  ·  FN: 0",
             hx + Inches(0.3), gy + Inches(2.45),
             panel_w - Inches(0.6), Inches(0.3),
             size=12, bold=True, color=INK, font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(s,
             "The one false positive (non_compete) flags a section that is really 'Category Exclusivity'. "
             "Honest failure mode worth discussing in viva — a real limit of the catalog on creator contracts.",
             hx + Inches(0.3), gy + Inches(2.8),
             panel_w - Inches(0.6), Inches(0.8),
             size=10, color=INK_SOFT, font=BODY_FONT, line_spacing=1.35)

    # tests strip at bottom
    add_round(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
              fill=CREAM, line=LINE)
    add_text(s, "21 / 21 unit tests pass · pipeline, exports, Q&A, score spread, held-out evaluation, calibration, redline rendering.",
             Inches(0.9), Inches(6.58), Inches(11.6), Inches(0.45),
             size=11, bold=True, color=FOREST_DEEP, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(s, page, total)


def slide_spread_calibration(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Score spread + calibration", "The model is well-calibrated and the gauge has real range.")

    # left: score spread bar chart
    lx = Inches(0.7)
    lw = Inches(6.0)
    lh = Inches(5.35)
    add_round(s, lx, Inches(1.7), lw, lh, fill=WHITE, line=LINE)
    add_text(s, "Risk score across the 5 bundled demos",
             lx + Inches(0.3), Inches(1.85),
             lw - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK, font=HEADING_FONT)
    add_text(s, "Same pipeline, no per-contract tuning. Each demo lands in a distinct band.",
             lx + Inches(0.3), Inches(2.25),
             lw - Inches(0.5), Inches(0.4),
             size=11, color=MUTED, font=BODY_FONT)

    demos = [
        ("High-risk freelance", 90, CORAL_DEEP, "Critical"),
        ("SaaS subscription", 85, CORAL, "Critical"),
        ("Held-out influencer", 80, CORAL, "Critical"),
        ("Balanced agreement", 53, AMBER, "Review"),
        ("Friendly letter", 25, GREEN, "Low Risk"),
    ]
    bar_max_w = Inches(3.5)
    bar_x = lx + Inches(2.05)
    for i, (label, val, color, badge) in enumerate(demos):
        y = Inches(2.85) + i * Inches(0.85)
        add_text(s, label, lx + Inches(0.3), y,
                 Inches(1.75), Inches(0.4),
                 size=11, bold=True, color=INK, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, bar_x, y + Inches(0.05), bar_max_w * (val/100), Inches(0.3),
                 fill=color)
        # value at end of bar
        add_text(s, str(val), bar_x + bar_max_w * (val/100) + Inches(0.1), y + Inches(0.05),
                 Inches(0.5), Inches(0.3),
                 size=11, bold=True, color=color, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
        # badge
        add_round(s, bar_x, y + Inches(0.4), Inches(0.95), Inches(0.25),
                  fill=PAPER, radius=0.45)
        add_text(s, badge, bar_x, y + Inches(0.4), Inches(0.95), Inches(0.25),
                 size=8, bold=True, color=color, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # right: calibration
    rx = Inches(7.0)
    rw = Inches(5.65)
    add_round(s, rx, Inches(1.7), rw, lh, fill=WHITE, line=LINE)
    add_text(s, "Confidence calibration",
             rx + Inches(0.3), Inches(1.85),
             rw - Inches(0.5), Inches(0.4),
             size=14, bold=True, color=INK, font=HEADING_FONT)
    add_text(s, "Mean confidence vs. empirical precision per bin, across 29 gold + held-out findings.",
             rx + Inches(0.3), Inches(2.25),
             rw - Inches(0.5), Inches(0.5),
             size=11, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    # calibration bins as grouped bar
    bins = [
        ("[0.55–0.70)", 0.66, 0.83, 6),
        ("[0.70–0.85)", 0.77, 1.00, 8),
        ("[0.85–1.00)", 0.95, 1.00, 15),
    ]
    bin_y = Inches(3.0)
    bar_h = Inches(0.32)
    bin_label_w = Inches(1.5)
    bin_bar_max = Inches(2.5)
    for i, (label, mean_conf, precision, n) in enumerate(bins):
        block_y = bin_y + i * Inches(1.0)
        add_text(s, label, rx + Inches(0.3), block_y,
                 bin_label_w, Inches(0.3),
                 size=10, bold=True, color=INK, font=BODY_FONT)
        add_text(s, f"n = {n}", rx + Inches(0.3), block_y + Inches(0.3),
                 bin_label_w, Inches(0.3),
                 size=9, color=MUTED, font=BODY_FONT)
        # mean confidence bar
        add_rect(s, rx + Inches(1.85), block_y, bin_bar_max * mean_conf, bar_h,
                 fill=MOSS)
        add_text(s, f"conf {mean_conf:.2f}",
                 rx + Inches(1.85) + bin_bar_max * mean_conf + Inches(0.1), block_y,
                 Inches(1.2), bar_h,
                 size=9, color=INK, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
        # precision bar
        add_rect(s, rx + Inches(1.85), block_y + Inches(0.4), bin_bar_max * precision, bar_h,
                 fill=FOREST_BRIGHT)
        add_text(s, f"prec {precision:.2f}",
                 rx + Inches(1.85) + bin_bar_max * precision + Inches(0.1), block_y + Inches(0.4),
                 Inches(1.2), bar_h,
                 size=9, color=INK, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # metrics card
    add_round(s, rx + Inches(0.3), Inches(6.05), rw - Inches(0.6), Inches(0.85),
              fill=CREAM)
    add_text(s, "Brier = 0.053  ·  ECE = 0.126  ·  overall precision = 0.97",
             rx + Inches(0.3), Inches(6.1),
             rw - Inches(0.6), Inches(0.35),
             size=11, bold=True, color=FOREST_DEEP, font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(s, "Model is slightly under-confident: when it says 0.95 it is right 100% of the time.",
             rx + Inches(0.3), Inches(6.4),
             rw - Inches(0.6), Inches(0.4),
             size=10, italic=True, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    slide_footer(s, page, total)


def slide_stack_limits(prs, page, total):
    s = blank(prs)
    slide_background(s)
    slide_header(s, "Stack, limits, future work", "What's in the box — and what we'd build next.")

    # 3 columns: stack, limitations, future
    col_w = Inches(4.1)
    col_h = Inches(5.0)
    col_gap = Inches(0.18)
    cols_total_w = col_w * 3 + col_gap * 2
    start_x = (SLIDE_W - cols_total_w) / 2
    col_y = Inches(1.8)

    def col(left, eyebrow, title, color, bullets):
        add_round(s, left, col_y, col_w, col_h, fill=WHITE, line=LINE)
        add_rect(s, left, col_y, col_w, Inches(0.12), fill=color)
        add_text(s, eyebrow, left + Inches(0.3), col_y + Inches(0.25),
                 col_w - Inches(0.5), Inches(0.3),
                 size=10, bold=True, color=MUTED, font=BODY_FONT)
        add_text(s, title, left + Inches(0.3), col_y + Inches(0.55),
                 col_w - Inches(0.5), Inches(0.55),
                 size=17, bold=True, color=INK, font=HEADING_FONT)
        add_bullets(s, bullets, left + Inches(0.3), col_y + Inches(1.2),
                    col_w - Inches(0.5), col_h - Inches(1.3),
                    size=11, color=INK_SOFT, bullet_color=color, line_spacing=1.45)

    col(start_x, "01 / TECH STACK", "Reproducible, lightweight", FOREST_BRIGHT, [
        "Streamlit · interactive UI with custom CSS",
        "scikit-learn · TF-IDF + cosine similarity",
        "Plotly · gauge, donut, radar, calibration chart",
        "pdfplumber + python-docx · document parsing",
        "Groq · optional LLM polish (works offline by default)",
    ])

    col(start_x + (col_w + col_gap), "02 / KNOWN LIMITS", "Honest weak spots", CORAL, [
        "Hand-curated catalog → recall capped to known clauses",
        "English only (no Urdu / Arabic support yet)",
        "Section splitter assumes paragraph-block formatting",
        "False positives possible when fair text mentions risky terms",
        "No version-comparison or diff view",
    ])

    col(start_x + (col_w + col_gap) * 2, "03 / WHAT'S NEXT", "If we kept building", AMBER, [
        "Sentence-Transformer classifier trained on CUAD",
        "Inline contract diff between two versions",
        "Multi-language support (Urdu + Arabic clauses)",
        "Real PDF export with WeasyPrint",
        "Cloud deployment + share-link mode for clients",
    ])

    slide_footer(s, page, total)


def slide_thank_you(prs, page, total):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=FOREST_DEEP)
    add_rect(s, 0, 0, SLIDE_W, Inches(3.8), fill=FOREST)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.4), fill=FOREST_BRIGHT)
    add_rect(s, 0, Inches(1.4), SLIDE_W, Inches(0.06), fill=GOLD)

    add_text(s, "LIVE DEMO · QUESTIONS",
             0, Inches(0.6), SLIDE_W, Inches(0.5),
             size=12, bold=True, color=CREAM, font=BODY_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Thank you", 0, Inches(2.3), SLIDE_W, Inches(1.4),
             size=72, bold=True, color=WHITE, font=HEADING_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s,
             "Open the Streamlit app and walk the examiner through the high-risk demo,",
             0, Inches(4.0), SLIDE_W, Inches(0.4),
             size=14, color=CREAM, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s,
             "then switch to the friendly demo to show the gauge in green.",
             0, Inches(4.35), SLIDE_W, Inches(0.4),
             size=14, color=CREAM, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # repo card
    add_round(s, Inches(2.5), Inches(5.1), Inches(8.3), Inches(1.7), fill=WHITE)
    add_text(s, "REPOSITORY", Inches(2.7), Inches(5.25), Inches(8), Inches(0.3),
             size=10, bold=True, color=MUTED, font=BODY_FONT)
    add_text(s, "github.com/ShehrozRafaqat/ClearClause-NLP-Project",
             Inches(2.7), Inches(5.55), Inches(8), Inches(0.55),
             size=20, bold=True, color=FOREST_DEEP, font=HEADING_FONT)
    add_text(s, "Shehroz Ali (MSDSF25M012) · Arslan Ahmad (MSDSF25M001)",
             Inches(2.7), Inches(6.15), Inches(8), Inches(0.4),
             size=12, color=MUTED, font=BODY_FONT)
    add_text(s, "Streamlit · scikit-learn · Plotly · optional Groq",
             Inches(2.7), Inches(6.45), Inches(8), Inches(0.3),
             size=11, italic=True, color=FOREST_BRIGHT, font=BODY_FONT)


# -----------------------------------------------------------------------------
# Build
# -----------------------------------------------------------------------------

def build_presentation(output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_builders = [
        slide_title,
        slide_problem,
        slide_solution,
        slide_pipeline,
        slide_catalog,
        slide_hybrid_extraction,
        slide_scoring,
        slide_negotiation,
        slide_redline,
        slide_qa,
        slide_dashboard,
        slide_evaluation,
        slide_spread_calibration,
        slide_stack_limits,
        slide_thank_you,
    ]
    total = len(slide_builders)
    for idx, builder in enumerate(slide_builders, 1):
        builder(prs, idx, total)

    prs.save(str(output_path))
    return output_path


def main() -> None:
    here = Path(__file__).resolve().parent
    output = here.parent / "ClearClause_Presentation.pptx"
    path = build_presentation(output)
    print(f"Saved: {path}  ({path.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
